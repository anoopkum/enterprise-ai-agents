"""
Parsers for born-digital formats. Text is already machine-readable here, so we use
free local libraries (not Document Intelligence) — cheaper and, for spreadsheets,
gives structured rows we can turn into graph nodes.
"""
import logging
import os

from src.ingestion.models import ParsedDocument
from src.ingestion.doc_classifier import classify

logger = logging.getLogger(__name__)


def parse_docx(path: str) -> ParsedDocument:
    from docx import Document

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    return ParsedDocument(path, "docx", text, doc_type=classify(os.path.basename(path), text))


def parse_pptx(path: str) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
    text = "\n".join(parts)
    return ParsedDocument(path, "pptx", text, doc_type=classify(os.path.basename(path), text))


def parse_html(path: str) -> ParsedDocument:
    from bs4 import BeautifulSoup

    with open(path, encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n").strip()
    return ParsedDocument(path, "html", text, doc_type=classify(os.path.basename(path), text))


def parse_markdown(path: str) -> ParsedDocument:
    with open(path, encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return ParsedDocument(path, "md", text, doc_type=classify(os.path.basename(path), text))


def parse_text(path: str) -> ParsedDocument:
    with open(path, encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return ParsedDocument(path, "txt", text, doc_type=classify(os.path.basename(path), text))


def parse_excel(path: str) -> ParsedDocument:
    """
    .xlsx/.xls → one text block per sheet (header row + data rows serialised).
    Structured rows are preserved in metadata['rows'] for the graph builder.
    """
    import pandas as pd

    engine = "xlrd" if path.lower().endswith(".xls") else "openpyxl"
    try:
        sheets = pd.read_excel(path, sheet_name=None, engine=engine, dtype=str)
    except Exception as exc:
        logger.warning("Excel parse failed for %s with %s: %s", path, engine, exc)
        sheets = {}

    parts: list[str] = []
    all_rows: list[dict] = []
    for sheet_name, df in sheets.items():
        df = df.fillna("")
        parts.append(f"# Sheet: {sheet_name}")
        parts.append(" | ".join(str(c) for c in df.columns))
        for _, row in df.iterrows():
            record = {str(k): str(v) for k, v in row.to_dict().items()}
            all_rows.append({"_sheet": sheet_name, **record})
            parts.append(" | ".join(str(v) for v in row.tolist()))

    text = "\n".join(parts)
    return ParsedDocument(
        path, "xls" if engine == "xlrd" else "xlsx", text,
        doc_type=classify(os.path.basename(path), text),
        metadata={"rows": all_rows, "sheet_count": len(sheets)},
    )
